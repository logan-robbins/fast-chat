"""
Image conversion utilities for PDF and PPTX files.

This module converts document pages/slides to PNG images for vision model
processing. Uses pure Python libraries - no LibreOffice required.

Dependencies:
- pdf2image>=1.17.0: PDF to image conversion (requires poppler-utils system package)
  Uses convert_from_bytes() with dpi and page range parameters.
- python-pptx>=0.6.23: Direct PPTX slide extraction
- Pillow>=10.2.0: Image manipulation and base64 encoding

System Requirements:
- poppler-utils: Required by pdf2image for PDF rendering
  Install: apt-get install poppler-utils (Debian/Ubuntu)
           brew install poppler (macOS)

SDK Versions (verified 02/05/2026):
- pdf2image 1.17.0: convert_from_bytes() with dpi, first_page, last_page params
- Pillow 10.2.0: Image.Resampling.LANCZOS for high-quality downscaling

Last Grunted: 02/05/2026
"""
from __future__ import annotations

import base64
import logging
from dataclasses import dataclass
from io import BytesIO

from pdf2image import convert_from_bytes
from PIL import Image

logger = logging.getLogger(__name__)


@dataclass
class DocumentImage:
    """
    Represents a converted document page or slide as a base64-encoded image.
    
    Immutable dataclass for passing page images through the extraction pipeline.
    The base64 data is suitable for direct use in OpenAI vision API requests.
    
    Attributes:
        index (int): 1-based page/slide number (matches PDF page numbering)
        base64_data (str): PNG image data encoded as base64 string (no data URI prefix).
            To use with OpenAI API, prepend "data:image/png;base64,"
    
    Example:
        image_url = f"data:image/png;base64,{doc_image.base64_data}"
    
    Last Grunted: 02/05/2026
    """
    index: int
    base64_data: str


# Default configuration -- matches IngestionConfig defaults.
# 200 DPI + 1568px max edge gives best quality at GPT-4o's native detail tier.
DEFAULT_DPI = 200
DEFAULT_MAX_EDGE = 1568
DEFAULT_MAX_PAGES = 20


def pdf_bytes_to_base64_images(
    pdf_bytes: bytes,
    *,
    max_pages: int | None = DEFAULT_MAX_PAGES,
    dpi: int = DEFAULT_DPI,
    max_edge: int = DEFAULT_MAX_EDGE,
) -> list[DocumentImage]:
    """
    Convert PDF document to list of base64-encoded PNG images.
    
    Uses pdf2image.convert_from_bytes() which wraps poppler's pdftoppm utility.
    Images are rendered at specified DPI then resized to fit within max_edge
    for optimal GPT-4o vision processing.
    
    Args:
        pdf_bytes (bytes): Raw PDF file bytes (any valid PDF)
        max_pages (int, optional): Maximum pages to convert (default: 20).
            Pages beyond this limit are skipped silently.
        dpi (int): Render resolution in dots per inch (default: 150).
            Higher DPI = better quality but larger images.
        max_edge (int): Maximum dimension after resizing (default: 1400px).
            Optimal for GPT-4o vision which handles up to 2048px.
    
    Returns:
        list[DocumentImage]: List of DocumentImage objects with 1-based indices.
            Empty list if conversion fails (logged as exception).
    
    Raises:
        No exceptions raised - failures return empty list and log error.
    
    Note:
        Requires poppler-utils system package. Without it, pdf2image raises
        PDFInfoNotInstalledError.
    
    Last Grunted: 02/05/2026
    """
    try:
        pages = convert_from_bytes(
            pdf_bytes,
            dpi=dpi,
            first_page=1,
            last_page=max_pages,
        )
        
        images = [
            DocumentImage(
                index=i,
                base64_data=_pil_to_base64(page, max_edge=max_edge),
            )
            for i, page in enumerate(pages, 1)
        ]
        
        logger.info(f"Converted PDF: {len(images)} pages")
        return images
        
    except Exception as exc:
        logger.exception(f"PDF conversion failed: {exc}")
        return []


def pptx_bytes_to_base64_images(
    pptx_bytes: bytes,
    *,
    max_pages: int | None = DEFAULT_MAX_PAGES,  # Named max_pages for API consistency
    dpi: int = DEFAULT_DPI,
    max_edge: int = DEFAULT_MAX_EDGE,
) -> list[DocumentImage]:
    """
    Convert PPTX presentation to list of base64-encoded PNG images.
    
    Uses python-pptx for direct slide extraction - no LibreOffice needed.
    Renders slides to images using Pillow with a simplified renderer that
    extracts text and basic shapes. Complex elements may not render perfectly
    but GPT-4o vision can interpret the content.
    
    Args:
        pptx_bytes (bytes): Raw PPTX file bytes (Office Open XML format)
        max_pages (int, optional): Maximum slides to convert (default: 20).
            Named max_pages for API consistency with pdf_bytes_to_base64_images.
        dpi (int): Used for EMU to pixel conversion (default: 150).
            Formula: pixels = emu * dpi / 914400 (914400 EMU = 1 inch)
        max_edge (int): Maximum image dimension after resizing (default: 1400px)
    
    Returns:
        list[DocumentImage]: List of DocumentImage objects with 1-based indices.
            Failed slides get placeholder images with error message.
            Empty list if PPTX parsing fails entirely.
    
    Raises:
        No exceptions raised - failures return empty list or placeholder images.
    
    Note:
        - Requires python-pptx>=0.6.23
        - Text rendering quality depends on available system fonts
        - Tables are rendered with simplified cell borders
    
    Last Grunted: 02/05/2026
    """
    try:
        from pptx import Presentation
        
        buffer = BytesIO(pptx_bytes)
        prs = Presentation(buffer)
        
        # Get slide dimensions
        slide_width = prs.slide_width.emu
        slide_height = prs.slide_height.emu
        
        # Calculate image dimensions (EMUs to pixels at roughly 96 DPI equivalent)
        # 914400 EMUs = 1 inch, so at 150 DPI: pixels = emu * 150 / 914400
        scale = dpi / 914400
        img_width = int(slide_width * scale)
        img_height = int(slide_height * scale)
        
        # Ensure reasonable dimensions
        img_width = max(800, min(img_width, 1920))
        img_height = max(600, min(img_height, 1080))
        
        images = []
        slides_to_process = list(prs.slides)[:max_pages] if max_pages else prs.slides
        
        for i, slide in enumerate(slides_to_process, 1):
            try:
                # Render slide to image
                slide_image = _render_slide_to_image(slide, img_width, img_height)
                if slide_image:
                    images.append(DocumentImage(
                        index=i,
                        base64_data=_pil_to_base64(slide_image, max_edge=max_edge),
                    ))
            except Exception as slide_exc:
                logger.warning(f"Failed to render slide {i}: {slide_exc}")
                # Create placeholder for failed slide
                placeholder = _create_placeholder_image(
                    img_width, img_height, 
                    f"Slide {i}\n(render failed)"
                )
                images.append(DocumentImage(
                    index=i,
                    base64_data=_pil_to_base64(placeholder, max_edge=max_edge),
                ))
        
        logger.info(f"Converted PPTX: {len(images)} slides")
        return images
        
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        return []
    except Exception as exc:
        logger.exception(f"PPTX conversion failed: {exc}")
        return []


def _render_slide_to_image(slide, width: int, height: int) -> Image.Image | None:
    """
    Render a PPTX slide to a PIL Image.
    
    Creates a white background and renders text/shapes from the slide.
    This is a simplified renderer - complex elements may not render perfectly,
    but the vision model can interpret the content.
    """
    from PIL import ImageDraw, ImageFont

    # Create white background
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)

    # Try to use a basic font, fall back to default
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
    except (OSError, IOError):
        font = ImageFont.load_default()
        title_font = font

    for shape in slide.shapes:
        try:
            # Get shape position (convert EMUs to pixels)
            if not hasattr(shape, 'left') or shape.left is None:
                continue
                
            x = int(shape.left.emu * width / slide.part.slide_layout.slide_master.slide_width.emu)
            y = int(shape.top.emu * height / slide.part.slide_layout.slide_master.slide_height.emu)
            
            # Render text frames
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    # Check if this might be a title (larger, near top)
                    is_title = y < height * 0.2 and hasattr(shape, 'is_placeholder')
                    use_font = title_font if is_title else font
                    
                    # Word wrap text
                    max_width = int(shape.width.emu * width / slide.part.slide_layout.slide_master.slide_width.emu)
                    wrapped_text = _wrap_text(text, use_font, max_width, draw)
                    draw.multiline_text((x, y), wrapped_text, fill='black', font=use_font)
            
            # Render tables
            if shape.has_table:
                table = shape.table
                cell_height = 20
                cell_y = y
                for row in table.rows:
                    cell_x = x
                    for cell in row.cells:
                        cell_text = cell.text[:50]  # Truncate long cells
                        draw.rectangle([cell_x, cell_y, cell_x + 100, cell_y + cell_height], outline='gray')
                        draw.text((cell_x + 2, cell_y + 2), cell_text, fill='black', font=font)
                        cell_x += 100
                    cell_y += cell_height
                    
        except Exception as shape_exc:
            # Skip shapes that fail to render
            continue
    
    return img


def _wrap_text(text: str, font, max_width: int, draw) -> str:
    """Wrap text to fit within max_width pixels."""
    words = text.split()
    lines = []
    current_line = []
    
    for word in words:
        test_line = ' '.join(current_line + [word])
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
            current_line = [word]
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return '\n'.join(lines)


def _create_placeholder_image(width: int, height: int, text: str) -> Image.Image:
    """Create a placeholder image with centered text."""
    from PIL import ImageDraw, ImageFont
    
    img = Image.new('RGB', (width, height), '#f0f0f0')
    draw = ImageDraw.Draw(img)
    
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
    except (OSError, IOError):
        font = ImageFont.load_default()
    
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    
    x = (width - text_width) // 2
    y = (height - text_height) // 2
    
    draw.text((x, y), text, fill='#666666', font=font)
    return img


def _pil_to_base64(image: Image.Image, *, max_edge: int) -> str:
    """
    Convert PIL image to base64-encoded PNG string.
    
    Resizes if either dimension exceeds max_edge, maintaining aspect ratio.
    Uses LANCZOS resampling for high-quality downscaling (best for photos
    and documents with fine detail).
    
    Args:
        image (Image.Image): PIL Image object (any mode - RGB, RGBA, L, etc.)
        max_edge (int): Maximum allowed dimension in pixels. If either width
            or height exceeds this, image is scaled down proportionally.
    
    Returns:
        str: Base64-encoded PNG data (without data URI prefix).
            To use with OpenAI API: f"data:image/png;base64,{result}"
    
    Formula:
        If max(width, height) > max_edge:
            ratio = max_edge / max(width, height)
            new_size = (int(width * ratio), int(height * ratio))
    
    Note:
        PNG format is used with optimize=True for smaller file size.
        Output is always RGB-compatible for vision API consumption.
    
    Last Grunted: 02/05/2026
    """
    if max(image.size) > max_edge:
        ratio = max_edge / max(image.size)
        new_size = (int(image.width * ratio), int(image.height * ratio))
        image = image.resize(new_size, Image.Resampling.LANCZOS)
    
    buffer = BytesIO()
    image.save(buffer, format="PNG", optimize=True)
    buffer.seek(0)
    return base64.b64encode(buffer.getvalue()).decode("utf-8")
